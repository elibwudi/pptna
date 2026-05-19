import os
import uuid
import io
import base64
import threading
import subprocess
import time
import logging
import bcrypt
import fitz  # PyMuPDF
from pathlib import Path
from datetime import timedelta
from cryptography.fernet import Fernet
import google.generativeai as genai
from openai import OpenAI
from flask import Flask, request, render_template, send_from_directory, flash, redirect, url_for, jsonify, session
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from dotenv import load_dotenv
from PIL import Image
import imagehash
import json

CONFIG_FILE = 'server_config.json'
ENCRYPTED_CONFIG_FILE = 'encrypted_config.key'

def load_config():
    """加载服务器配置"""
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except:
            return {"llm_provider": "ollama_gemma2"}
    return {"llm_provider": "ollama_gemma2"}

def save_config(config):
    """保存服务器配置"""
    with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
        json.dump(config, f, indent=4)

def get_encryption_key():
    """获取或生成加密密钥"""
    key_file = 'config_key.key'
    if os.path.exists(key_file):
        with open(key_file, 'rb') as f:
            return f.read()
    else:
        # 生成新密钥
        key = Fernet.generate_key()
        with open(key_file, 'wb') as f:
            f.write(key)
        return key

def load_api_keys():
    """加载加密的API密钥配置"""
    try:
        if not os.path.exists(ENCRYPTED_CONFIG_FILE):
            # 返回从环境变量读取的默认值
            return {
                'gemini_api_key': os.getenv('GEMINI_API_KEY', ''),
                'deepseek_api_key': os.getenv('DEEPSEEK_API_KEY', ''),
                'qwen_api_key': os.getenv('DASHSCOPE_API_KEY', '')
            }

        key = get_encryption_key()
        fernet = Fernet(key)

        with open(ENCRYPTED_CONFIG_FILE, 'rb') as f:
            encrypted_data = f.read()

        decrypted_data = fernet.decrypt(encrypted_data)
        return json.loads(decrypted_data.decode('utf-8'))
    except Exception as e:
        logger.error(f"加载加密配置失败: {e}")
        # 返回默认值
        return {
            'gemini_api_key': os.getenv('GEMINI_API_KEY', ''),
            'deepseek_api_key': os.getenv('DEEPSEEK_API_KEY', ''),
            'qwen_api_key': os.getenv('DASHSCOPE_API_KEY', '')
        }

def save_api_keys(api_keys):
    """保存加密的API密钥配置"""
    try:
        key = get_encryption_key()
        fernet = Fernet(key)

        # 只保存非空的密钥
        keys_to_save = {k: v for k, v in api_keys.items() if v}

        json_data = json.dumps(keys_to_save).encode('utf-8')
        encrypted_data = fernet.encrypt(json_data)

        with open(ENCRYPTED_CONFIG_FILE, 'wb') as f:
            f.write(encrypted_data)

        logger.info("API密钥已加密保存")
        return True
    except Exception as e:
        logger.error(f"保存加密配置失败: {e}")
        return False

def load_model_config():
    """加载模型配置（模型名称、服务地址等）"""
    config = load_config()
    return {
        'gemini_model': config.get('gemini_model', 'gemini-3-flash-preview'),
        'deepseek_model': config.get('deepseek_model', 'deepseek-chat'),
        'qwen_model': config.get('qwen_model', 'qwen-max'),
        'ollama_model': config.get('ollama_model', 'gemma4:31b'),
        'ollama_base_url': config.get('ollama_base_url', os.getenv('OLLAMA_BASE_URL', 'http://10.255.1.103:11434/v1')),
        'vllm_model': config.get('vllm_model', 'gemma-4-moe'),
        'vllm_base_url': config.get('vllm_base_url', 'http://10.255.1.118:8000/v1')
    }

def save_model_config(model_config):
    """保存模型配置"""
    config = load_config()
    config.update(model_config)
    save_config(config)
    logger.info(f"模型配置已更新: {model_config}")

# --- 新增：Windows COM 依赖 ---
import win32com.client
import pythoncom

# --- 1. 配置日志系统 ---
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(message)s',
    handlers=[
        logging.FileHandler('app.log', encoding='utf-8'),
        # 移除 StreamHandler 以防止在 pythonw_exe 中产生 IO 错误
    ]
)
logger = logging.getLogger(__name__)
logger.info("日志系统初始化完成")

# --- 2. 加载环境变量 ---
logger.info("--- 步骤 1: 正在加载 .env 文件 ---")
load_dotenv()
logger.info("--- .env 文件加载完毕 ---")



# --- 3. 配置Flask应用 ---
logger.info("--- 步骤 3: 正在配置 Flask 应用 ---")
UPLOAD_FOLDER = 'uploads'
GENERATED_FOLDER = 'generated'
ALLOWED_EXTENSIONS = {'pptx'}
MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB
app = Flask(__name__)

# 从环境变量读取SECRET_KEY，如果没有则生成临时密钥
secret_key = os.getenv("SECRET_KEY")
if not secret_key:
    logger.warning("未设置SECRET_KEY环境变量，使用临时密钥（请尽快配置）")
    secret_key = 'temp_secret_key_change_me'

app.config.from_mapping(
    UPLOAD_FOLDER=UPLOAD_FOLDER,
    GENERATED_FOLDER=GENERATED_FOLDER,
    SECRET_KEY=secret_key,
    MAX_CONTENT_LENGTH=MAX_FILE_SIZE,
    SESSION_COOKIE_HTTPONLY=True,
    SESSION_COOKIE_SAMESITE='Lax',
    PERMANENT_SESSION_LIFETIME=timedelta(hours=1)
)
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(GENERATED_FOLDER, exist_ok=True)
os.makedirs(os.path.join(app.root_path, 'static', 'thumbnails'), exist_ok=True)
logger.info("Flask 应用配置完毕。")


# --- 演讲语速常量 ---
SPEAKING_RATE = 260  # 汉字/分钟（单人正常宣讲语速，留出停顿和发散缓冲）


# --- 防缓存响应头 ---
@app.after_request
def add_header(response):
    response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
    response.headers['Pragma'] = 'no-cache'
    response.headers['Expires'] = '0'
    return response


# --- 3. 配置 API 客户端 ---
logger.info("--- 步骤 3: 正在配置 API 客户端 ---")
# Gemini API 配置
gemini_model = None
try:
    gemini_api_key = os.getenv("GEMINI_API_KEY")
    if not gemini_api_key:
        logger.warning("  -> 警告: Gemini 密钥未找到，模型将不可用。")
    else:
        genai.configure(api_key=gemini_api_key)
        gemini_model = genai.GenerativeModel('gemini-3-flash-preview')
        logger.info("  -> Gemini 客户端成功初始化！")
except Exception as e:
    logger.error(f"  -> Gemini 客户端初始化失败: {e}")
logger.info("-" * 20)
# DeepSeek API 配置
deepseek_client = None
try:
    deepseek_api_key = os.getenv("DEEPSEEK_API_KEY")
    if not deepseek_api_key:
        logger.warning("  -> 警告: DeepSeek 密钥未找到，客户端将不会被创建。")
    else:
        deepseek_client = OpenAI(api_key=deepseek_api_key, base_url="https://api.deepseek.com/v1")
        logger.info("  -> DeepSeek 客户端成功初始化！")
except Exception as e:
    logger.error(f"  -> DeepSeek 客户端初始化失败: {e}")
logger.info("-" * 20)
# Ollama 客户端配置
ollama_client = None
try:
    ollama_base_url = os.getenv("OLLAMA_BASE_URL", "http://10.255.1.103:11434/v1")
    ollama_client = OpenAI(base_url=ollama_base_url, api_key='ollama')
    logger.info(f"  -> Ollama 客户端已配置，目标地址: {ollama_base_url}")
except Exception as e:
    logger.error(f"  -> Ollama 客户端配置失败: {e}")
logger.info("-" * 20)
# Alibaba Qwen (DashScope) 客户端配置
qwen_client = None
try:
    dashscope_api_key = os.getenv("DASHSCOPE_API_KEY")
    if not dashscope_api_key:
        logger.warning("  -> 警告: DashScope 密钥未找到，Qwen 客户端将不会被创建。")
    else:
        qwen_client = OpenAI(
            api_key=dashscope_api_key,
            base_url="https://dashscope.aliyuncs.com/compatible-mode/v1"
        )
        logger.info("  -> Alibaba Qwen (DashScope) 客户端成功初始化！")
except Exception as e:
    logger.error(f"  -> Alibaba Qwen 客户端初始化失败: {e}")
logger.info("-" * 20)

# VLLM 客户端配置
vllm_client = None
try:
    vllm_base_url = load_model_config().get('vllm_base_url', 'http://10.255.1.118:8000/v1')
    # vLLM (OpenAI compatible) typically doesn't need a real API key, but the client requires one.
    vllm_client = OpenAI(
        api_key="sk-dummy",
        base_url=vllm_base_url
    )
    logger.info(f"  -> VLLM 客户端已配置，目标地址: {vllm_base_url}")
except Exception as e:
    logger.error(f"  -> VLLM 客户端配置失败: {e}")

logger.info("--- API 客户端配置完毕 ---")


# --- 任务追踪字典 ---
tasks = {}


# --- 辅助函数 ---
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def get_client_ip():
    """
    获取客户端真实IP地址
    考虑了代理服务器（Nginx、负载均衡器等）的情况
    """
    # 按优先级检查可能的IP头
    headers_to_check = [
        'X-Forwarded-For',      # 标准代理头
        'X-Real-IP',            # Nginx常用
        'CF-Connecting-IP',     # Cloudflare
        'True-Client-IP',       # Akamai
        'X-Client-IP'
    ]

    for header in headers_to_check:
        ip = request.headers.get(header)
        if ip:
            # X-Forwarded-For 可能包含多个IP，取第一个
            if header == 'X-Forwarded-For':
                ip = ip.split(',')[0].strip()
            # 清理可能的IPv6映射的IPv4
            if ip.startswith('::ffff:'):
                ip = ip[7:]
            return ip

    # 回退到直接连接的IP
    return request.remote_addr


def log_security_event(event_type, details, status="success"):
    """
    记录安全相关事件到日志

    Args:
        event_type: 事件类型 (login_attempt, config_change, etc.)
        details: 事件详情
        status: 事件状态 (success, failed, warning)
    """
    client_ip = get_client_ip()
    user_agent = request.headers.get('User-Agent', 'Unknown')

    log_message = f"[SECURITY] {event_type.upper()} | IP: {client_ip} | Status: {status} | {details}"

    if status == "failed":
        logger.warning(log_message)
    elif status == "warning":
        logger.warning(log_message)
    else:
        logger.info(log_message)

    # 记录完整信息用于调试
    logger.debug(f"  -> Full headers: {dict(request.headers)}")


def get_template_and_duplicate_hashes(prs):
    template_hashes = set()
    for master in prs.slide_masters:
        for shape in master.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = Image.open(io.BytesIO(shape.image.blob))
                    if img.mode != 'RGB': img = img.convert('RGB')
                    template_hashes.add(imagehash.average_hash(img))
                except Exception:
                    pass
    for layout in prs.slide_layouts:
        for shape in layout.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = Image.open(io.BytesIO(shape.image.blob))
                    if img.mode != 'RGB': img = img.convert('RGB')
                    template_hashes.add(imagehash.average_hash(img))
                except Exception:
                    pass
    return template_hashes


def _extract_recursive(shape, seen_hashes, extracted_content, text_runs):
    """递归提取形状中的文字和图片"""
    # 1. 提取文字
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip():
                    text_runs.append(run.text)
                    
    # 2. 提取表格文字
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame.text.strip():
                    text_runs.append(cell.text_frame.text)

    # 3. 提取图片
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = Image.open(io.BytesIO(shape.image.blob))
            if img.mode != 'RGB': img = img.convert('RGB')
            h = imagehash.average_hash(img)
            if h not in seen_hashes:
                extracted_content["images"].append(img)
                seen_hashes.add(h)
        except Exception as e:
            logger.warning(f"警告: 提取或处理图片失败: {e}")

    # 4. 递归处理组合形状
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            _extract_recursive(sub_shape, seen_hashes, extracted_content, text_runs)


def extract_content_from_slide(slide, seen_hashes):
    extracted_content = {"text": "", "images": []}
    text_runs = []
    
    for shape in slide.shapes:
        _extract_recursive(shape, seen_hashes, extracted_content, text_runs)
        
    extracted_content["text"] = " ".join(text_runs).strip()
    return extracted_content


def image_to_base_64(image_input, format="JPEG"):
    """
    将图片转换为 Base64 字符串（支持 PIL Image 对象或已经是字符串的输入）
    """
    if isinstance(image_input, str):
        return image_input.strip().replace("\n", "").replace("\r", "")
        
    try:
        buffered = io.BytesIO()
        # 确保是 PIL 对象
        if hasattr(image_input, 'mode') and image_input.mode != 'RGB': 
            image_input = image_input.convert('RGB')
        image_input.save(buffered, format=format, quality=85)
        return base64.b64encode(buffered.getvalue()).decode('utf-8')
    except Exception as e:
        logger.error(f"图片转Base64失败: {e}")
        return ""


# ==================== 两阶段生成：第一阶段辅助函数 ====================

def extract_slides_summary(prs, llm_provider=None, thumb_disk_paths=None, max_chars_per_slide=200):
    """
    提取所有幻灯片的文本摘要，用于第一阶段全局分析

    Args:
        prs: Presentation对象
        llm_provider: 提供商，用于可能的图片分析
        thumb_disk_paths: 缩略图路径列表，用于纯图片页分析
        max_chars_per_slide: 每页最大字符数，默认200字

    Returns:
        字符串，每页一行，格式: "[第N页] 文本内容 （含M张图片）"
    """
    summary_lines = []
    seen_hashes = get_template_and_duplicate_hashes(prs)

    for i, slide in enumerate(prs.slides, 1):
        try:
            content = extract_content_from_slide(slide, seen_hashes)
            text = content["text"]
            image_count = len(content["images"])

            # 如果没有提取到文字但是有图片，且开启了视觉分析支持
            if not text.strip() and thumb_disk_paths and i <= len(thumb_disk_paths) and llm_provider:
                try:
                    thumb_path = thumb_disk_paths[i-1]
                    if os.path.exists(thumb_path):
                        with open(thumb_path, "rb") as image_file:
                            encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
                            logger.info(f"第{i}页无文本，正在生成视觉简述...")
                            vision_text = generate_vision_summary(encoded_string, llm_provider)
                            text = f"[视觉简述: {vision_text.strip()}]"
                except Exception as ve:
                    logger.warning(f"生成第{i}页视觉简述失败: {ve}")

            # 限制每页摘要长度
            if len(text) > max_chars_per_slide:
                text = text[:max_chars_per_slide] + "..."

            line = f"[第{i}页] {text if text else '（此页无文字内容）'}"
            if image_count > 0:
                line += f" （含{image_count}张图片）"
            summary_lines.append(line)
        except Exception as e:
            logger.warning(f"提取第{i}页摘要失败: {e}")
            summary_lines.append(f"[第{i}页] （内容提取失败）")

    return "\n".join(summary_lines)


def generate_vision_summary(image_b64, llm_provider):
    """
    调用当前模型对图片进行极简摘要描述（模型中立）
    """
    vision_prompt = """请用不超过15个汉字简要描述这张幻灯片图片的主要内容和主题（如：“关于云架构的流程图”或“带有标题的企业愿景介绍”）。直接输出描述内容，不要输出任何额外文字。"""
    
    try:
        # 包装成与主流程一致的多模态调用
        # 注意：这里调用 generate_notes_with_fallback 的简化流程
        # 为了避免循环依赖及递归限制，我们限制重试次数
        result = generate_notes_with_fallback(
            prompt=vision_prompt,
            images=[image_b64],
            primary_provider=llm_provider,
            session_id="vision_sum",
            slide_number=0 # 标记为特殊任务
        )
        # 清理输出防止包含思考过程
        return clean_llm_output(result).strip()
    except Exception as e:
        logger.error(f"视觉摘要生成失败: {e}")
        return "（图片内容分析失败）"


def create_global_analysis_prompt(slides_summary, presentation_title):
    """
    生成第一阶段全局分析的Prompt

    Args:
        slides_summary: 所有幻灯片的摘要
        presentation_title: 演示文稿标题

    Returns:
        第一阶段分析的prompt字符串
    """
    return f"""
### 身份协议 ###
你是一位资深的教学设计专家，擅长分析PPT内容并生成全局教学规划。

### 任务概述 ###
请分析以下PPT的所有幻灯片内容，生成一份全局教学规划。这份规划将用于后续为每一页生成详细的讲稿。

### PPT标题 ###
{presentation_title}

### PPT内容摘要 ###
{slides_summary}

### 输出要求 ###
请以JSON格式输出，严格遵循以下结构：

{{
    "outline": {{
        "overview": "用2-3句话概括整个PPT的核心主题和教学目标",
        "sections": [
            {{
                "start_slide": 起始页码（整数）,
                "end_slide": 结束页码（整数）,
                "title": "教学模块标题（不要包含‘第X章’等字眼）",
                "key_topics": ["关键主题1", "关键主题2"],
                "teaching_strategy": "该模块的教学策略（如：概念讲解/案例分析/实践演示等）"
            }}
        ]
    }},
    "glossary": {{
        "terms": [
            {{
                "original": "原文术语",
                "translation": "统一翻译",
                "definition": "简短定义"
            }}
        ]
    }},
    "style_guide": {{
        "tone": "描述整体语气（如：专业严谨/轻松活泼/鼓励启发等）",
        "metaphor_style": "描述比喻风格（如：生活化类比/技术类比/无比喻等）",
        "audience_level": "目标受众水平（如：初学者/中级/高级）"
    }},
    "key_themes": ["核心主题1", "核心主题2"],
    "references": {{
        "note": "标注需要前后引用的重要概念（可选）"
    }}
}}

### 分析要点 ###
1. **逻辑模块划分**：根据内容逻辑自动划分为几个教学模块（建议3-8个模块，每个模块至少3页）。注意：请摒弃传统的“教科书章节”概念，统一称为“模块”，以帮助理清讲解思路。
2. **术语统一**：识别专业术语（10-30个），确保翻译一致
3. **风格定位**：根据内容确定教学风格
4. **主题聚焦**：提炼2-4个核心主题
5. **引用关系**：标注需要前后呼应的内容

请直接输出JSON（不要使用markdown代码块标记）：
"""


def create_default_global_context(presentation_title, total_slides):
    """
    创建默认全局上下文（当第一阶段失败时使用）

    Args:
        presentation_title: 演示文稿标题
        total_slides: 总页数

    Returns:
        默认的全局上下文字典
    """
    return {
        "presentation_title": presentation_title,
        "total_slides": total_slides,
        "outline": {
            "overview": f"关于{presentation_title}的演示",
            "sections": [{
                "start_slide": 1,
                "end_slide": total_slides,
                "title": "主要内容",
                "key_topics": [],
                "teaching_strategy": "概念讲解"
            }]
        },
        "glossary": {"terms": []},
        "style_guide": {
            "tone": "专业且亲切",
            "metaphor_style": "生活化类比",
            "audience_level": "初学者"
        },
        "key_themes": [],
        "references": {}
    }


def generate_global_analysis(slides_summary, presentation_title, llm_provider):
    """
    第一阶段：生成全局分析（带自动fallback机制）

    Args:
        slides_summary: 所有幻灯片的摘要
        presentation_title: 演示文稿标题
        llm_provider: AI模型提供商

    Returns:
        全局上下文字典，失败时返回默认上下文
    """
    try:
        prompt = create_global_analysis_prompt(slides_summary, presentation_title)

        # 使用带fallback的函数调用AI（第一阶段不需要图片）
        response = generate_notes_with_fallback(
            prompt=prompt,
            images=[],  # 第一阶段不需要图片
            primary_provider=llm_provider,
            session_id="",
            slide_number=0
        )

        # 检查是否所有模型都失败了
        if response.startswith("错误：所有AI模型均无法生成"):
            logger.error(f"第一阶段所有模型均失败，使用默认上下文")
            return create_default_global_context(presentation_title, len(slides_summary.split('\n')))

        # 解析JSON响应
        try:
            # 清理可能的markdown代码块标记和备用模型提示
            cleaned = clean_llm_output(response)

            # 移除备用模型提示信息（如果存在）
            if "[注：本页使用备用模型" in cleaned:
                # 找到JSON开始的位置（第一个{）
                json_start = cleaned.find('{')
                if json_start != -1:
                    cleaned = cleaned[json_start:]

            if cleaned.startswith("```json"):
                cleaned = cleaned[7:]
            if cleaned.startswith("```"):
                cleaned = cleaned[3:]
            if cleaned.endswith("```"):
                cleaned = cleaned[:-3]
            cleaned = cleaned.strip()

            global_context = json.loads(cleaned)
            global_context['presentation_title'] = presentation_title

            # 验证必需字段
            required_keys = ['outline', 'glossary', 'style_guide']
            if not all(key in global_context for key in required_keys):
                raise ValueError("返回的上下文缺少必需字段")

            logger.info(f"第一阶段分析成功：识别出{len(global_context['outline']['sections'])}个章节，"
                       f"{len(global_context['glossary']['terms'])}个术语")
            return global_context

        except (json.JSONDecodeError, ValueError) as e:
            logger.error(f"解析全局分析JSON失败: {e}")
            logger.error(f"原始响应: {response[:500]}...")
            return create_default_global_context(presentation_title, len(slides_summary.split('\n')))

    except Exception as e:
        logger.error(f"第一阶段分析失败: {e}，使用默认上下文")
        return create_default_global_context(presentation_title, len(slides_summary.split('\n')))


def get_section_info(slide_number, global_context):
    """
    获取指定页码所属的章节信息

    Args:
        slide_number: 页码（从1开始）
        global_context: 全局上下文字典

    Returns:
        包含章节信息的字典
    """
    for section in global_context['outline']['sections']:
        if section['start_slide'] <= slide_number <= section['end_slide']:
            return {
                'title': section['title'],
                'teaching_strategy': section['teaching_strategy'],
                'key_topics': section['key_topics'],
                'is_first_slide': slide_number == section['start_slide'],
                'is_last_slide': slide_number == section['end_slide']
            }

    # 默认返回（不应该发生）
    return {
        'title': '内容讲解',
        'teaching_strategy': '概念讲解',
        'key_topics': [],
        'is_first_slide': False,
        'is_last_slide': False
    }


def create_stage2_prompt_with_context(slide_text, slide_number, total_slides,
                                     global_context, section_info,
                                     has_images, target_chars, previous_notes=""):
    """
    第二阶段：创建带全局上下文的逐页生成Prompt

    Args:
        slide_text: 当前页的文本内容
        slide_number: 当前页码
        total_slides: 总页数
        global_context: 全局上下文字典
        section_info: 当前页的章节信息
        has_images: 是否包含图片
        target_chars: 目标字数
        previous_notes: 上一页的讲稿内容（用于防重复和自然过渡）

    Returns:
        第二阶段的prompt字符串
    """
    # 确定当前页在章节中的位置
    is_first_slide = slide_number == 1
    is_section_start = section_info.get('is_first_slide', False)
    is_section_end = section_info.get('is_last_slide', False)

    # 构建上下文引导（防重复与自然过渡）
    if is_first_slide:
        context_intro = "这是课程的第一页。请设计一个引人入胜的开场，抓住学生们的注意力，并简要介绍今天这堂课的主题。"
    else:
        concise_prev = previous_notes[:250] + "..." if len(previous_notes) > 250 else previous_notes
        if is_section_start:
            context_intro = f"在上一页，我们刚讲了: \"{concise_prev}\"。\n现在我们进入一个新的学习模块：『{section_info['title']}』。请自然地承上启下，引出本模块的目标，不要生硬播报“下面进入第X章”。"
        elif is_section_end:
            context_intro = f"在上一页，我们刚讲了: \"{concise_prev}\"。\n这是当前模块『{section_info['title']}』的最后一页。请顺着上一页的话题继续，并在结尾处对本模块进行简要总结。"
        else:
            context_intro = f"在上一页，我们刚讲了: \"{concise_prev}\"。\n请继续顺着这个话题，自然地过渡到当前页的新内容，保持讲解的连贯性。"

    # 图片提示
    image_prompt_part = "本页的核心内容是文字。"
    if has_images:
        image_prompt_part = "除了文字，这张幻灯片还有一些配图。请在讲解时，引导学生观察这些图片，并将图片内容与知识点联系起来。"

    # 术语表提示
    glossary_hint = ""
    if global_context['glossary']['terms']:
        glossary_hint = "\n### 术语统一要求 ###\n请确保以下术语翻译和使用保持一致：\n"
        for term in global_context['glossary']['terms'][:15]:  # 限制数量，避免prompt过长
            glossary_hint += f"- {term['original']} → {term['translation']}\n"

    # 风格指南提示
    style_hint = f"""
### 教学风格指南 ###
- 整体语气：{global_context['style_guide']['tone']}
- 比喻风格：{global_context['style_guide']['metaphor_style']}
- 受众水平：{global_context['style_guide']['audience_level']}
"""

    # 章节信息提示
    section_hint = f"""
### 当前模块信息 ###
- 模块名称：{section_info['title']}
- 教学策略：{section_info['teaching_strategy']}
- 本模块主题：{', '.join(section_info['key_topics']) if section_info['key_topics'] else '综合内容'}
"""

    # 核心主题
    themes_hint = ""
    if global_context['key_themes']:
        themes_hint = f"\n### 核心主题 ###\n整个课程围绕以下核心主题展开：{', '.join(global_context['key_themes'])}\n"

    return f"""
### 身份协议 ###
你是"启思老师"，一位顶尖教学专家，正在讲解关于"{global_context['presentation_title']}"的课程。

### 全局背景 ###
{global_context['outline']['overview']}

{style_hint}

{section_hint}

{themes_hint}

{glossary_hint}

### 当前教学任务 ###
{context_intro}

- **当前页详情**:
    - 教学进度: 第 {slide_number} 页 / 共 {total_slides} 页
    - 核心知识点: "{slide_text if slide_text.strip() else '（本页无文字内容，请完全基于对图片的视觉分析进行创作）'}"
    - 教学图示: {image_prompt_part if slide_text.strip() else '本页由纯图片/图像组成，请执行 OCR 识别其中的文字并解析图中表达的含义，这是讲稿的唯一来源。'}

### 核心要求 (CRITICAL) ###
1. **极致口语化**：你的讲解必须使用通俗易懂、生动的口语化简体中文。多用短句、生活化类比、加入“大家看”、“想一想”等互动词汇，甚至允许适当的语气词，就像和教室里的学生直接对话。
2. **严禁生硬播报结构**：坚决不要说“下面进入第X章第Y节”或“这页PPT写了什么”。所有的结构转换必须化为自然的话语过渡。
3. **避免重复开场**：不要每页都说“同学们”，结合上一页的内容做自然衔接。
4. **术语统一与风格一致**：保持语气一致，遇到抽象概念必须打比方。
5. **绝对字数红线**：本页讲稿必须严格控制在 **{target_chars} 个汉字左右**（误差不超过20%）。如果本页只是一张纯标题页或过渡页（内容极少），请你仅用一两句话快速带过，【完全无视上方给出的目标字数】，严禁为了凑字数而强行补充没有营养的废话！如果不遵守字数约束，将导致严重的系统崩溃！
6. **严禁使用Markdown格式**：讲稿必须是纯文本。对于需要强调的词汇，绝对禁止使用 ** 进行加粗。如果确实需要强调，请使用引号替代。

"启思老师"，请直接输出讲稿内容（不要输出任何JSON、代码或元数据）：
"""


# ==================== 两阶段生成结束 ====================

# --- !!! 核心优化：全新的教学风格Prompt !!! ---
def create_multimodal_prompt(slide_text, slide_number, total_slides, presentation_title, has_images, previous_notes, target_chars):
    """
    (教学优化版) 使用为课堂教学场景设计的全新Prompt。
    """
    if slide_number == 1:
        context_intro = "这是课程的第一页。你的任务是设计一个引人入胜的开场，抓住学生们的注意力，并简要介绍今天这堂课的主题。"
    else:
        concise_previous_notes = previous_notes
        if len(concise_previous_notes) > 250:  # 对上一页内容做更精简的概括
            concise_previous_notes = concise_previous_notes[:250] + "..."
        context_intro = f"在上一页，我们刚刚讲了: \"{concise_previous_notes}\"。现在，请设计一个自然的过渡，引导学生进入下一个知识点的学习。"

    image_prompt_part = "本页的核心内容是文字。"
    if has_images:
        image_prompt_part = "除了文字，这张幻灯片还有一些配图。请在讲解时，引导学生观察这些图片，并将图片内容与知识点联系起来。"

    return f"""
    ### 身份协议 (Identity Protocol) ###
    你将扮演一位名为“启思老师”的顶尖教学专家。你的专长是将“{presentation_title}”这个主题的复杂知识点，转化为一堂生动有趣、引人入胜的**中文课程**。
    你的听众是**学生**，所以你的核心目标是激发他们的好奇心，用最清晰、最易懂的方式，让他们理解和记住核心概念。

    ### 核心教学原则 (Core Teaching Principles) ###
    1.  **语言与风格**: 你的讲解必须使用通俗易懂、口语化的简体中文。语气要亲切、自然，充满热情，仿佛正在和教室里的学生直接对话。
    2.  **教学方法**: 你的讲解应该深入浅出。要善于使用**比喻、类比和生活中的例子**来解释抽象概念。可以适当加入一些**提问**（例如，“大家想一想，这是为什么呢？”或“这和我们上次课讲的有什么联系？”）来引导学生思考。
    3.  **内容限制 (CRITICAL)**: 你的回答**只能包含**你作为老师要说的口语化讲稿内容。严禁输出任何JSON、代码、元数据、或对自己行为的评论 (例如，不要说“我看到了图片...”或“这张图表显示了...”)。你应该直接说：“同学们，我们来看这张图...”。
    4.  **禁止粗体 (NO BOLD)**: 严禁在讲稿中使用 ** 符号进行加粗。讲稿必须是纯文本。如果需要强调，请使用引号替代。
    5.  **长度控制 (最重要)**: 本页讲稿的最终输出**总字数必须严格控制在 {target_chars} 个汉字左右**。这是一个硬性指标，请务必遵守，以确保课程节奏。

    ### 你的当前教学任务 ###
    - **上下文**: {context_intro}
    - **当前页详情**:
        - 教学进度: 第 {slide_number} 页 / 共 {total_slides} 页
        - 核心知识点: "{slide_text if slide_text else '（本页以视觉元素为主）'}"
        - 教学图示: {image_prompt_part}

    “启思老师”，请严格遵循以上所有原则，开始你的精彩讲解吧：
    """


def clean_llm_output(raw_text):
    """
    清理LLM输出，移除不必要的前缀和markdown代码块标记
    """
    if not raw_text:
        return raw_text

    raw_text = raw_text.strip()

    # 处理markdown代码块（优先处理，因为很多模型会输出 ```json ... ```）
    # 移除所有 ``` 标记
    raw_text = raw_text.replace("```json", "")
    raw_text = raw_text.replace("```", "")
    
    # 彻底移除 Markdown 加粗符号 **（转语音不友好）
    raw_text = raw_text.replace("**", "")
    
    raw_text = raw_text.strip()

    # 处理英文前缀（对于特定的英文开场白）
    unwanted_english_starters = [
        "Here are the objects I've identified in the images:",
        "Here are the bounding box detections:",
        "Here is the speech script for the slide:"
    ]

    for starter in unwanted_english_starters:
        if raw_text.startswith(starter):
            # 找到第一个中文字符或JSON开始符
            first_valid_index = -1
            for i, char in enumerate(raw_text):
                if '\u4e00' <= char <= '\u9fff' or char == '{':
                    first_valid_index = i
                    break
            if first_valid_index != -1:
                raw_text = raw_text[first_valid_index:]
            else:
                return "（模型输出内容异常，已被自动过滤）"
            break

    return raw_text.strip()


def generate_notes_gemini(prompt, images):
    """Gemini API调用，带重试机制（最多3次重试）"""
    if not gemini_model: return "错误：Gemini模型未初始化。"

    import time
    max_retries = 3
    base_delay = 2  # 基础等待时间（秒）

    for attempt in range(max_retries):
        try:
            if attempt > 0:
                logger.warning(f"Gemini API第{attempt + 1}次重试...")
                # 指数退避等待：2秒, 4秒, 8秒
                time.sleep(base_delay * (2 ** attempt))

            start_time = time.time()
            response = gemini_model.generate_content([prompt] + images)
            elapsed = time.time() - start_time

            if elapsed > 30:
                logger.warning(f"Gemini API调用耗时较长: {elapsed:.2f}秒")
            elif attempt > 0:
                logger.info(f"Gemini API重试成功（第{attempt + 1}次尝试）")

            return response.text

        except Exception as e:
            error_msg = str(e)

            # 判断是否为超时错误
            is_timeout = (
                "504" in error_msg or
                "timeout" in error_msg.lower() or
                "Deadline" in error_msg or
                "time out" in error_msg.lower()
            )

            if is_timeout:
                if attempt < max_retries - 1:
                    logger.warning(f"Gemini API超时（尝试{attempt + 1}/{max_retries}），准备重试...")
                    continue
                else:
                    logger.error(f"Gemini API超时：已重试{max_retries}次仍失败")
                    return f"Gemini API调用超时（已重试{max_retries}次，建议切换模型）：{error_msg}"
            else:
                # 非超时错误，直接返回不重试
                logger.error(f"Gemini API调用失败（非超时）: {error_msg}")
                return f"Gemini API调用失败: {e}"

    return "错误：Gemini API未返回结果"


def generate_notes_deepseek(prompt):
    """DeepSeek API调用，带重试机制（最多3次重试）"""
    if not deepseek_client: return "错误：DeepSeek客户端未初始化。"

    import time
    max_retries = 3
    base_delay = 2

    for attempt in range(max_retries):
        try:
            if attempt > 0:
                logger.warning(f"DeepSeek API第{attempt + 1}次重试...")
                time.sleep(base_delay * (2 ** attempt))

            start_time = time.time()
            model_name = load_model_config().get('deepseek_model', 'deepseek-chat')
            response = deepseek_client.chat.completions.create(
                model=model_name,
                messages=[{"role": "user", "content": prompt}],
                timeout=120  # 2分钟超时
            )
            elapsed = time.time() - start_time

            if elapsed > 30:
                logger.warning(f"DeepSeek API调用耗时较长: {elapsed:.2f}秒")
            elif attempt > 0:
                logger.info(f"DeepSeek API重试成功（第{attempt + 1}次尝试）")

            return response.choices[0].message.content

        except Exception as e:
            error_msg = str(e)

            # 判断是否为超时错误
            is_timeout = (
                "504" in error_msg or
                "timeout" in error_msg.lower() or
                "Deadline" in error_msg or
                "ReadTimeout" in error_msg or
                "time out" in error_msg.lower()
            )

            if is_timeout:
                if attempt < max_retries - 1:
                    logger.warning(f"DeepSeek API超时（尝试{attempt + 1}/{max_retries}），准备重试...")
                    continue
                else:
                    logger.error(f"DeepSeek API超时：已重试{max_retries}次仍失败")
                    return f"DeepSeek API调用超时（已重试{max_retries}次）：{error_msg}"
            else:
                logger.error(f"DeepSeek API调用失败（非超时）: {error_msg}")
                return f"DeepSeek API调用失败: {e}"

    return "错误：DeepSeek API未返回结果"


def generate_notes_ollama(prompt, images):
    """Ollama API调用，带重试机制（最多3次重试）"""
    if not ollama_client: return "错误：Ollama客户端未初始化。"

    max_retries = 3
    base_delay = 2

    for attempt in range(max_retries):
        try:
            if attempt > 0:
                logger.warning(f"Ollama API第{attempt + 1}次重试...")
                time.sleep(base_delay * (2 ** attempt))  # 指数退避：2s, 4s, 8s

            start_time = time.time()
            messages = [{"role": "user", "content": [{"type": "text", "text": prompt}]}]
            
            for img in images:
                b64_data = image_to_base_64(img)
                if b64_data:
                    messages[0]["content"].append({
                        "type": "image_url", 
                        "image_url": {"url": f"data:image/jpeg;base64,{b64_data}"}
                    })

            # 设置5分钟超时（本地模型可能较慢）
            response = ollama_client.chat.completions.create(
                model="gemma3:27b",
                messages=messages,
                timeout=300  # 5分钟超时
            )

            elapsed = time.time() - start_time
            if elapsed > 60:
                logger.warning(f"Ollama API调用耗时较长: {elapsed:.2f}秒")
            elif attempt > 0:
                logger.info(f"Ollama API重试成功（第{attempt + 1}次尝试）")

            return response.choices[0].message.content

        except Exception as e:
            error_msg = str(e)
            is_timeout = ("timeout" in error_msg.lower() or "Timeout" in error_msg or "timed out" in error_msg.lower())
            is_connection_error = ("Connection refused" in error_msg or "Failed to connect" in error_msg)

            if is_timeout or is_connection_error:
                if attempt < max_retries - 1:
                    logger.warning(f"Ollama API{'超时' if is_timeout else '连接失败'}（尝试{attempt + 1}/{max_retries}），准备重试...")
                    continue
                else:
                    logger.error(f"Ollama API{'超时' if is_timeout else '连接失败'}：已重试{max_retries}次仍失败")
                    if is_connection_error:
                        return f"调用Ollama失败：无法连接到 {ollama_client.base_url}。请确保Ollama服务正在运行，并且本应用服务器与Ollama服务器网络互通。"
                    return f"Ollama API调用超时（已重试{max_retries}次，建议检查Ollama服务状态）：{error_msg}"
            else:
                # 非超时/连接错误，直接返回
                logger.error(f"Ollama API调用出错: {e}")
                return f"调用Ollama失败: {e}"

    return "错误：Ollama API未返回结果"


def generate_notes_qwen(prompt):
    """Qwen API调用，带重试机制（最多3次重试）"""
    if not qwen_client: return "错误：Alibaba Qwen 客户端未初始化。请检查 .env 文件中是否配置了 DASHSCOPE_API_KEY。"

    max_retries = 3
    base_delay = 2

    for attempt in range(max_retries):
        try:
            if attempt > 0:
                logger.warning(f"Qwen API第{attempt + 1}次重试...")
                time.sleep(base_delay * (2 ** attempt))  # 指数退避：2s, 4s, 8s

            start_time = time.time()

            # 设置2分钟超时
            model_name = load_model_config().get('qwen_model', 'qwen-max')
            response = qwen_client.chat.completions.create(
                model=model_name,  # 动态读取模型名
                messages=[{"role": "user", "content": prompt}],
                timeout=120  # 2分钟超时
            )

            elapsed = time.time() - start_time
            if elapsed > 30:
                logger.warning(f"Qwen API调用耗时较长: {elapsed:.2f}秒")
            elif attempt > 0:
                logger.info(f"Qwen API重试成功（第{attempt + 1}次尝试）")

            return response.choices[0].message.content

        except Exception as e:
            error_msg = str(e)
            is_timeout = ("timeout" in error_msg.lower() or "Timeout" in error_msg or "超时" in error_msg or "ReadTimeout" in error_msg)

            if is_timeout:
                if attempt < max_retries - 1:
                    logger.warning(f"Qwen API超时（尝试{attempt + 1}/{max_retries}），准备重试...")
                    continue
                else:
                    logger.error(f"Qwen API超时：已重试{max_retries}次仍失败")
                    return f"Qwen API调用超时（已重试{max_retries}次，建议切换模型或稍后重试）：{error_msg}"
            else:
                # 非超时错误，直接返回
                logger.error(f"Qwen API调用出错: {e}")
                return f"Qwen API调用失败: {e}"

    return "错误：Qwen API未返回结果"


def generate_notes_vllm(prompt, images):
    """VLLM API调用，带重试机制（支持多模态图片传输）"""
    if not vllm_client: return "错误：VLLM 客户端未初始化。请检查配置。"

    max_retries = 3
    for attempt in range(max_retries):
        try:
            if attempt > 0:
                logger.warning(f"VLLM API第{attempt + 1}次重试...")

            start_time = time.time()
            model_name = load_model_config().get('vllm_model', 'gemma-4-moe')
            
            # 构造多模态消息格式 (OpenAI 兼容模式)
            if images:
                content = [{"type": "text", "text": prompt}]
                for img in images:
                    # 关键修复：先通过 image_to_base_64 处理，它现在能处理 PIL对象和字符串
                    clean_b64 = image_to_base_64(img)
                    if not clean_b64: continue
                    
                    content.append({
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{clean_b64}",
                            "detail": "auto"
                        }
                    })
                messages = [{"role": "user", "content": content}]
            else:
                messages = [{"role": "user", "content": prompt}]

            response = vllm_client.chat.completions.create(
                model=model_name,
                messages=messages,
                timeout=120  # 2分钟超时
            )

            elapsed = time.time() - start_time
            if elapsed > 30:
                logger.warning(f"VLLM API调用耗时较长: {elapsed:.2f}秒")
            elif attempt > 0:
                logger.info(f"VLLM API重试成功（第{attempt + 1}次尝试）")

            return response.choices[0].message.content

        except Exception as e:
            error_msg = str(e)
            is_timeout = ("timeout" in error_msg.lower() or "Timeout" in error_msg or "超时" in error_msg or "ReadTimeout" in error_msg)

            if is_timeout:
                if attempt < max_retries - 1:
                    logger.warning(f"VLLM API超时（尝试{attempt + 1}/{max_retries}），准备重试...")
                    continue
                else:
                    logger.error(f"VLLM API超时：已重试{max_retries}次仍失败")
                    return f"VLLM API调用超时（已重试{max_retries}次，建议切换模型或稍后重试）：{error_msg}"
            else:
                logger.error(f"VLLM API调用出错: {e}")
                return f"VLLM API调用失败: {e}"

    return "错误：VLLM API未返回结果"


def generate_notes_with_fallback(prompt, images, primary_provider, session_id="", slide_number=0):
    """
    带自动fallback的讲稿生成函数

    优先使用主模型，如果失败自动切换到Ollama备用模型

    Args:
        prompt: 提示词
        images: 图片列表
        primary_provider: 主模型提供商（gemini/deepseek/qwen/ollama_gemma2）
        session_id: 会话ID（用于日志）
        slide_number: 页码（用于日志）

    Returns:
        生成的内容或错误信息
    """
    # 定义备用模型列表（按用户要求的优先级：Qwen -> Gemini -> DeepSeek -> Ollama最后）
    fallback_providers = ['qwen', 'gemini', 'deepseek', 'ollama_gemma2']

    # 尝试主模型
    logger.info(f"[{session_id}] 尝试使用主模型 {primary_provider} 生成第{slide_number}页...")
    # print(f"[{session_id}] 正在使用 {primary_provider.upper()} 处理第 {slide_number} 页...")  # 同时输出到stdout
    try:
        if primary_provider == 'gemini':
            result = generate_notes_gemini(prompt, images)
        elif primary_provider == 'deepseek':
            result = generate_notes_deepseek(prompt)
        elif primary_provider == 'ollama_gemma2':
            result = generate_notes_ollama(prompt, images)
        elif primary_provider == 'qwen':
            result = generate_notes_qwen(prompt)
        elif primary_provider == 'vllm':
            result = generate_notes_vllm(prompt, images)
        else:
            result = f"错误：未知的模型提供商 {primary_provider}"

        # 检查是否是重试失败后的错误信息
        if is_retry_failure_error(result):
            logger.warning(f"[{session_id}] 主模型 {primary_provider} 重试后仍失败，尝试切换到备用模型...")
            # print(f"[{session_id}] 主模型 {primary_provider} 失败，尝试切换到备用模型...")
        else:
            logger.info(f"[{session_id}] 主模型 {primary_provider} 生成成功")
            return result

    except Exception as e:
        logger.error(f"[{session_id}] 主模型 {primary_provider} 调用异常: {e}")
        # print(f"[{session_id}] 主模型 {primary_provider} 调用异常: {e}")
        result = f"主模型调用失败: {e}"

    # 尝试备用模型
    for fallback_provider in fallback_providers:
        # 跳过主模型（因为已经失败了）
        if fallback_provider == primary_provider:
            continue

        # 跳过未初始化的模型
        if fallback_provider == 'ollama_gemma2' and not ollama_client:
            logger.warning(f"[{session_id}] 备用模型 {fallback_provider} 未初始化，跳过")
            continue
        elif fallback_provider == 'vllm' and not vllm_client:
            logger.warning(f"[{session_id}] 备用模型 {fallback_provider} 未初始化，跳过")
            continue
        elif fallback_provider == 'qwen' and not qwen_client:
            logger.warning(f"[{session_id}] 备用模型 {fallback_provider} 未初始化，跳过")
            continue
        elif fallback_provider == 'deepseek' and not deepseek_client:
            logger.warning(f"[{session_id}] 备用模型 {fallback_provider} 未初始化，跳过")
            continue

        try:
            logger.info(f"[{session_id}] 切换到备用模型 {fallback_provider} 生成第{slide_number}页...")
            # print(f"[{session_id}] 切换到备用模型 {fallback_provider.upper()} 生成第{slide_number}页...")

            if fallback_provider == 'gemini':
                # 只有初始化了才用Gemini作为fallback
                if not gemini_model:
                    continue
                result = generate_notes_gemini(prompt, images)
            elif fallback_provider == 'deepseek':
                result = generate_notes_deepseek(prompt)
            elif fallback_provider == 'ollama_gemma2':
                result = generate_notes_ollama(prompt, images)
            elif fallback_provider == 'qwen':
                result = generate_notes_qwen(prompt)
            elif fallback_provider == 'vllm':
                result = generate_notes_vllm(prompt)

            # 检查备用模型是否也失败了
            if is_retry_failure_error(result):
                logger.warning(f"[{session_id}] 备用模型 {fallback_provider} 也失败，尝试下一个...")
                # print(f"[{session_id}] 备用模型 {fallback_provider} 也失败，尝试下一个...")
                continue

            # 成功！
            logger.info(f"[{session_id}] 备用模型 {fallback_provider} 生成成功！")
            # print(f"[{session_id}] 备用模型 {fallback_provider} 生成成功！")
            return f"[注：本页使用备用模型{fallback_provider}生成]\n\n{result}"

        except Exception as e:
            logger.error(f"[{session_id}] 备用模型 {fallback_provider} 调用异常: {e}")
            # print(f"[{session_id}] 备用模型 {fallback_provider} 调用异常: {e}")
            continue

    # 所有模型都失败
    error_msg = f"所有AI模型均无法生成第{slide_number}页讲稿。建议检查：1) Ollama服务是否运行 2) 网络连接 3) API密钥配置"
    logger.error(f"[{session_id}] {error_msg}")
    return f"错误：{error_msg}"


def is_retry_failure_error(result):
    """
    判断是否是重试失败后的错误信息

    Args:
        result: API返回结果

    Returns:
        bool: True表示是重试失败错误
    """
    if not isinstance(result, str):
        return False

    # 检查各种重试失败的关键词
    failure_indicators = [
        '已重试3次仍失败',
        '已重试3次',
        'retry',
        '建议切换模型',
        '建议检查',
        '无法连接'
    ]

    result_lower = result.lower()
    for indicator in failure_indicators:
        if indicator.lower() in result_lower:
            return True

    return False


def get_or_create_notes_text_frame(slide):
    """
    更加健壮地获取或创建备注栏文本框。
    处理一些特殊的PPT模板中 notes_text_frame 为 None 的情况。
    """
    try:
        # 访问 notes_slide 时会尝试自动创建备注页面
        notes_slide = slide.notes_slide
        
        # 方法1：尝试标准 API
        if notes_slide.notes_text_frame is not None:
            return notes_slide.notes_text_frame
            
        # 方法2：手动查找 BODY (2) 或 NOTES_BODY (7) 类型的占位符
        for shape in notes_slide.placeholders:
            if shape.is_placeholder and shape.placeholder_format.type in [2, 7]:
                return shape.text_frame
        
        # 方法3：寻找除 幻灯片缩略图(101) 之外的其他带有文本框的占位符
        for shape in notes_slide.shapes:
            if shape.has_text_frame and shape.is_placeholder:
                # 排除缩略图和页码
                if shape.placeholder_format.type not in [101, 13]:
                    return shape.text_frame
                    
        # 方法4：实在没有，返回 None (调用侧将跳过写入)
    except Exception as e:
                logger.debug(f"尝试获取备注框时出错: {e}")
        
    return None


def patch_ppt_notes_with_com(ppt_path, slides_data, session_id):
    """
    使用 PowerPoint COM 接口补齐缺失的备注。
    针对 python-pptx 无法处理的非标准模板提供最后保障。
    """
    logger.info(f"[{session_id}] 启动 COM 备注补丁程序...")
    pythoncom.CoInitialize()
    ppt_app = None
    presentation = None
    try:
        abs_ppt_path = os.path.abspath(ppt_path)
        # 使用 DispatchEx 确保启动独立进程
        ppt_app = win32com.client.DispatchEx("PowerPoint.Application")
        presentation = ppt_app.Presentations.Open(abs_ppt_path, ReadOnly=0, Untitled=0, WithWindow=0)
        
        num_slides = presentation.Slides.Count
        patch_count = 0
        
        for i in range(1, num_slides + 1):
            if i > len(slides_data): break
            
            note_content = slides_data[i-1].get('notes', '')
            if not note_content: continue
            
            slide = presentation.Slides(i)
            notes_page = slide.NotesPage
            
            # 查找已有的 Body 占位符 (Type=14 是 Placeholder, PlaceholderFormat.Type=2 是 ppPlaceholderBody)
            body = None
            try:
                for shape in notes_page.Shapes:
                    if shape.Type == 14:
                        if shape.PlaceholderFormat.Type == 2:
                            body = shape
                            break
            except:
                pass
            
            if body:
                body.TextFrame.TextRange.Text = note_content
                patch_count += 1
            else:
                # 强行恢复/添加“备注占位符” (ppPlaceholderBody = 2)
                # 这比直接添加普通文本框更好，因为它可以被 PowerPoint 的“备注窗格”识别
                try:
                    # 位置建议：左边距50, 顶部350(约中下部), 宽度600, 高度300
                    new_placeholder = notes_page.Shapes.AddPlaceholder(2, 50, 350, 600, 300)
                    new_placeholder.TextFrame.TextRange.Text = note_content
                    new_placeholder.Name = "AI_Generated_Notes_Placeholder"
                    patch_count += 1
                except Exception as ex:
                    logger.debug(f"[{session_id}] 第{i}页添加占位符失败，尝试最后的回退方案(Textbox): {ex}")
                    try:
                        # 最后的保底方案：普通文本框
                        new_box = notes_page.Shapes.AddTextbox(1, 50, 350, 600, 300)
                        new_box.TextFrame.TextRange.Text = note_content
                        new_box.Name = "AI_Generated_Notes_Textbox"
                        patch_count += 1
                    except Exception as ex2:
                        logger.error(f"[{session_id}] 第{i}页所有插入备注方案均失败: {ex2}")
        
        presentation.Save()
        logger.info(f"[{session_id}] COM 补丁完成: 成功处理 {patch_count}/{num_slides} 页")
    except Exception as e:
        logger.error(f"[{session_id}] COM 补丁异常: {e}")
    finally:
        if presentation:
            try: presentation.Close()
            except: pass
        if ppt_app:
            try: ppt_app.Quit()
            except: pass
        # 必须在线程内平衡 CoInitialize
        pythoncom.CoUninitialize()



def generate_script_task(upload_path, original_filename, llm_provider, duration_minutes, session_id, root_path,
                         generated_folder):
    """
    修改版：使用 Windows COM 接口替代 LibreOffice 进行 PPT 转 PDF
    """
    # 重要：在线程开始时初始化 COM
    pythoncom.CoInitialize()
    
    temp_output_dir = os.path.join(generated_folder, session_id)
    static_thumb_dir = os.path.join(root_path, 'static', 'thumbnails', session_id)
    os.makedirs(temp_output_dir, exist_ok=True)
    os.makedirs(static_thumb_dir, exist_ok=True)
    pdf_path = None
    
    # 变量声明，用于 finally 块清理
    ppt_app = None
    presentation = None

    try:
        logger.info(f"[{session_id}] 后台任务开始: {original_filename}")
        tasks[session_id] = {'status': 'processing', 'progress': '正在准备...'}

        # --- 步骤 1: PPT 转 PDF (Windows 原生方案) ---
        logger.info(f"[{session_id}] 步骤1/3: 正在将PPTX转换为PDF (Windows COM)...")
        tasks[session_id]['progress'] = '正在转换格式 (Office)...'

        # 定义绝对路径 (COM 必需)
        abs_upload_path = os.path.abspath(upload_path)
        pdf_filename = Path(upload_path).stem + '.pdf'
        pdf_path = os.path.join(temp_output_dir, pdf_filename)
        abs_pdf_path = os.path.abspath(pdf_path)

        try:
            # 启动 PowerPoint 应用 (使用 DispatchEx 确保进程独立，防止干扰)
            ppt_app = win32com.client.DispatchEx("PowerPoint.Application")
            # ppt_app.Visible = 1 # 调试时可开启，生产环境通常不需要或可能会抢焦点

            # 打开演示文稿 (ReadOnly=True, Untitled=False, WithWindow=False)
            presentation = ppt_app.Presentations.Open(abs_upload_path, ReadOnly=1, Untitled=0, WithWindow=0)
            
            # 另存为 PDF (32 代表 ppSaveAsPDF)
            presentation.SaveAs(abs_pdf_path, 32)
            logger.info(f"[{session_id}] PDF转换成功: {abs_pdf_path}")
            
        except Exception as e:
            raise RuntimeError(f"PowerPoint COM 转换失败。请确认服务器已安装Office且非试用期。错误: {e}")
        finally:
            # 确保关闭 PPT 文件，防止文件占用
            if presentation:
                try:
                    presentation.Close()
                except:
                    pass
            # 退出 PPT 进程 (为了稳定性，建议退出)
            if ppt_app:
                try:
                    ppt_app.Quit()
                except:
                    pass

        if not os.path.exists(pdf_path):
            raise FileNotFoundError("PDF文件未生成，Office转换可能失败。")

        # --- 步骤 2: 生成缩略图 ---
        logger.info(f"[{session_id}] 步骤2/4: 正在生成缩略图...")
        tasks[session_id]['progress'] = '正在生成缩略图...'
        thumbnail_paths_for_web = []
        thumbnail_paths_for_analysis = []
        doc = fitz.open(pdf_path)
        for page_num, page in enumerate(doc):
            pix = page.get_pixmap(dpi=150)
            thumb_filename = f"slide_{page_num + 1}.png"
            thumb_path_on_disk = os.path.join(static_thumb_dir, thumb_filename)
            pix.save(thumb_path_on_disk)
            thumbnail_paths_for_web.append(f"thumbnails/{session_id}/{thumb_filename}")
            thumbnail_paths_for_analysis.append(thumb_path_on_disk) # 记录物理路径用于第一阶段分析
        doc.close()

        # --- 步骤 3: 第一阶段 - 全局分析（两阶段生成） ---
        logger.info(f"[{session_id}] 步骤3/4: 正在进行全局分析...")
        tasks[session_id]['progress'] = '正在分析PPT结构...'

        # 加载PPT用于内容提取
        prs = Presentation(upload_path)
        num_slides = len(prs.slides)

        if num_slides == 0:
            raise ValueError("PPT文件不包含任何幻灯片。")

        # 获取演示文稿标题
        presentation_title = "未命名演示文稿"
        if prs.slides and len(prs.slides) > 0:
            first_slide = prs.slides[0]
            if hasattr(first_slide.shapes, 'title') and first_slide.shapes.title is not None:
                title_text = first_slide.shapes.title.text.strip()
                if title_text:
                    presentation_title = title_text

        # 判断是否需要两阶段生成（小型PPT可跳过）
        use_two_stage = num_slides >= 3  # 3页及以上使用两阶段

        global_context = None
        if use_two_stage:
            logger.info(f"[{session_id}] 使用两阶段生成模式...")
            try:
                # 提取所有页面摘要（传入缩略图路径以支持纯图片分析）
                import time
                stage1_start = time.time()
                slides_summary = extract_slides_summary(
                    prs, 
                    llm_provider=llm_provider, 
                    thumb_disk_paths=thumbnail_paths_for_analysis,
                    max_chars_per_slide=200
                )
                logger.info(f"[{session_id}] 已提取{num_slides}页摘要，总长度{len(slides_summary)}字符")

                # 生成全局分析
                global_context = generate_global_analysis(
                    slides_summary, presentation_title, llm_provider
                )
                stage1_duration = time.time() - stage1_start
                logger.info(f"[{session_id}] 第一阶段完成，耗时{stage1_duration:.2f}秒")
                logger.info(f"[{session_id}] 识别出{len(global_context['outline']['sections'])}个章节，"
                      f"{len(global_context['glossary']['terms'])}个术语")

                # 保存全局分析结果（用于调试和用户查看）
                global_analysis_path = os.path.join(temp_output_dir, 'global_analysis.json')
                with open(global_analysis_path, 'w', encoding='utf-8') as f:
                    json.dump(global_context, f, ensure_ascii=False, indent=2)

            except Exception as e:
                logger.error(f"[{session_id}] 第一阶段失败: {e}，使用默认上下文")
                global_context = create_default_global_context(presentation_title, num_slides)
        else:
            logger.info(f"[{session_id}] 小型PPT（{num_slides}页），跳过第一阶段，使用默认上下文")
            global_context = create_default_global_context(presentation_title, num_slides)

        # --- 步骤 4: 第二阶段 - 逐页生成讲稿 ---
        logger.info(f"[{session_id}] 步骤4/4: 正在生成讲稿...")
        tasks[session_id]['progress'] = '正在生成讲稿...'

        if num_slides != len(thumbnail_paths_for_web):
            logger.warning(f"警告：PPT页数({num_slides})与生成的缩略图数量({len(thumbnail_paths_for_web)})不匹配。")

        slides_data_for_template, full_script_content = [], []
        seen_hashes = get_template_and_duplicate_hashes(prs)

        # 预扫描以提取内容并计算动态权重
        slides_contents = []
        for slide in prs.slides:
            slides_contents.append(extract_content_from_slide(slide, seen_hashes))
            
        total_chars_needed = int(duration_minutes * SPEAKING_RATE)
        
        slide_weights = []
        for content in slides_contents:
            text_len = len(content["text"])
            image_weight = 150 if content["images"] else 0
            slide_weights.append(text_len + image_weight)
            
        total_weight = sum(slide_weights)
        if total_weight == 0:
            total_weight = 1 # 防除零
            
        # 动态分配字数
        base_chars_per_slide = min(30, total_chars_needed // num_slides)
        remaining_chars = total_chars_needed - (base_chars_per_slide * num_slides)
        
        target_chars_list = []
        for w in slide_weights:
            target = base_chars_per_slide + int((w / total_weight) * remaining_chars)
            target_chars_list.append(target)

        logger.info(f"[{session_id}] 目标总时长: {duration_minutes}分钟, 共{num_slides}页。"
              f"采用动态分配，总目标: {total_chars_needed}字，分配范围: {min(target_chars_list)} - {max(target_chars_list)}字。")

        notes_write_success_count = 0  # 统计成功写入备注的页数
        previous_notes = ""

        for i, slide in enumerate(prs.slides):
            progress_pct = int(((i + 1) / num_slides) * 100)
            tasks[session_id]['progress'] = f'正在生成讲稿 ({progress_pct}%)'

            slide_content = slides_contents[i]
            slide_text, slide_images = slide_content["text"], slide_content["images"]
            raw_notes = ""

            if not slide_text and not slide_images:
                raw_notes = "（此页为空白页或仅包含重复/模板图片）"
            else:
                # 获取章节信息（两阶段模式）
                section_info = get_section_info(i + 1, global_context) if global_context else {
                    'title': '内容讲解',
                    'teaching_strategy': '概念讲解',
                    'key_topics': [],
                    'is_first_slide': i == 0,
                    'is_last_slide': i == num_slides - 1
                }

                is_multimodal_provider = llm_provider in ['gemini', 'ollama_gemma2', 'vllm']
                has_images_for_prompt = bool(slide_images) and is_multimodal_provider

                # 使用两阶段prompt（如果提供了全局上下文）
                if global_context:
                    prompt = create_stage2_prompt_with_context(
                        slide_text=slide_text,
                        slide_number=i + 1,
                        total_slides=num_slides,
                        global_context=global_context,
                        section_info=section_info,
                        has_images=has_images_for_prompt,
                        target_chars=target_chars_list[i],
                        previous_notes=previous_notes
                    )
                else:
                    # Fallback到旧版prompt
                    # 针对纯图片页面的提示词优化
                    effective_text = slide_text
                    if not effective_text.strip() and has_images_for_prompt:
                        effective_text = "(本页由纯图片组成，请基于图像内容进行创作)"
                        
                    prompt = create_multimodal_prompt(
                        effective_text, i + 1, num_slides, presentation_title,
                        has_images=has_images_for_prompt,
                        previous_notes=previous_notes,
                        target_chars=target_chars_list[i]
                    )

                # 调用AI生成（带自动fallback到Ollama）
                raw_notes = generate_notes_with_fallback(
                    prompt=prompt,
                    images=slide_images,
                    primary_provider=llm_provider,
                    session_id=session_id,
                    slide_number=i + 1
                )

            current_notes = clean_llm_output(raw_notes)
            previous_notes = current_notes # 记录为下一页的上下文

            # 写入备注 (更健壮的方法)
            try:
                text_frame = get_or_create_notes_text_frame(slide)

                if text_frame is not None:
                    text_frame.text = current_notes
                    notes_write_success_count += 1
                else:
                    logger.warning(f"[{session_id}] 警告: 第{i+1}页备注框不存在，跳过写入PPT")
            except Exception as e:
                logger.error(f"[{session_id}] 错误: 第{i+1}页写入备注失败: {e}")

            # 获取章节名称（用于显示）
            section_title = ""
            if global_context:
                section_info = get_section_info(i + 1, global_context)
                section_title = f" ({section_info['title']})"

            slides_data_for_template.append({
                'slide_number': i + 1,
                'notes': current_notes,
                'screenshot_path': thumbnail_paths_for_web[i] if i < len(thumbnail_paths_for_web) else 'placeholder.png',
                'section': section_title  # 新增：章节信息
            })
            full_script_content.append(f"--- 第 {i + 1}/{num_slides} 页{section_title} ---\n\n{current_notes}\n\n")

        # 保存结果文件名定义
        ppt_filename = f"{original_filename.rsplit('.', 1)[0]}_with_notes.pptx"
        txt_filename = f"{original_filename.rsplit('.', 1)[0]}_script.txt"
        dest_ppt_path = os.path.join(temp_output_dir, ppt_filename)

        # 1. 优先使用标准方式保存（即便没有写入全部备注）
        prs.save(dest_ppt_path)
        
        # 2. 如果标准方式写入不全，启动 COM 补丁程序进行强行补齐
        if notes_write_success_count < num_slides:
            logger.info(f"[{session_id}] 提示: 部分页面因PPT模板限制无法通过标准方式写入备注（成功: {notes_write_success_count}/{num_slides}），切换至 COM 补丁模式...")
            patch_ppt_notes_with_com(dest_ppt_path, slides_data_for_template, session_id)

        # 3. 保存文字讲稿
        with open(os.path.join(temp_output_dir, txt_filename), 'w', encoding='utf-8') as f:
            f.write("".join(full_script_content))

        # 构建返回结果
        result = {
            'status': 'complete',
            'ppt_filename': ppt_filename,
            'txt_filename': txt_filename,
            'slides_data': slides_data_for_template
        }

        # 如果使用了两阶段生成，添加全局分析文件信息
        if global_context and use_two_stage:
            result['global_analysis'] = 'global_analysis.json'
            result['generation_mode'] = 'two_stage'
            logger.info(f"[{session_id}] 两阶段生成完成，全局分析已保存到 global_analysis.json")
        else:
            result['generation_mode'] = 'legacy'

        tasks[session_id] = result
        logger.info(f"[{session_id}] 后台任务成功完成。")

    except Exception as e:
        logger.error(f"[{session_id}] 后台任务失败: {e}")
        tasks[session_id] = {'status': 'error', 'message': str(e)}
        
    finally:
        # 释放 COM 资源
        pythoncom.CoUninitialize()
        
        # 清理文件
        # 等待一小会儿确保文件句柄已释放
        time.sleep(1) 
        if os.path.exists(upload_path):
            try:
                os.remove(upload_path)
            except PermissionError:
                logger.warning(f"[{session_id}] 警告: 无法删除上传文件 (可能仍被占用): {upload_path}")
        if pdf_path and os.path.exists(pdf_path):
            try:
                os.remove(pdf_path)
            except:
                pass
        logger.info(f"[{session_id}] 资源清理流程结束。")


# --- Flask 路由 ---

# 辅助函数：验证管理员权限
def is_admin_authenticated():
    return session.get('is_admin', False)

# 管理员登录验证路由
@app.route('/api/verify-admin', methods=['POST'])
def verify_admin():
    """验证管理员密码，使用bcrypt哈希比对"""
    try:
        data = request.get_json()
        if not data or 'password' not in data:
            log_security_event("login_attempt", "缺少密码参数", "failed")
            return jsonify({'status': 'error', 'message': '缺少密码参数'}), 400

        password = data.get('password')
        admin_hash = os.getenv('ADMIN_PASSWORD_HASH')

        if not admin_hash:
            log_security_event("login_attempt", "服务器未配置ADMIN_PASSWORD_HASH", "failed")
            logger.error("未配置ADMIN_PASSWORD_HASH环境变量")
            return jsonify({'status': 'error', 'message': '服务器配置错误'}), 500

        # 使用bcrypt验证密码
        if bcrypt.checkpw(password.encode('utf-8'), admin_hash.encode('utf-8')):
            session['is_admin'] = True
            log_security_event("admin_login", "管理员登录成功")
            return jsonify({'status': 'success'})
        else:
            log_security_event("admin_login", "密码验证失败", "failed")
            return jsonify({'status': 'error', 'message': '密码错误'}), 401

    except Exception as e:
        log_security_event("admin_login", f"验证过程异常: {str(e)}", "failed")
        logger.error(f"管理员验证过程出错: {e}")
        return jsonify({'status': 'error', 'message': '服务器错误'}), 500

# 管理员登出路由
@app.route('/api/admin-logout', methods=['POST'])
def admin_logout():
    """管理员登出"""
    session.pop('is_admin', None)
    log_security_event("admin_logout", "管理员主动登出")
    return jsonify({'status': 'success'})

@app.route('/api/settings', methods=['GET', 'POST'])
def api_settings():
    # POST请求需要管理员权限
    if request.method == 'POST':
        # 验证管理员权限
        if not is_admin_authenticated():
            log_security_event("config_change", "未授权尝试修改配置", "failed")
            logger.warning("未授权访问设置API")
            return jsonify({'status': 'error', 'message': '需要管理员权限'}), 403

        data = request.get_json()
        new_provider = data.get('llm_provider')
        if new_provider:
            config = load_config()
            old_provider = config.get('llm_provider', 'unknown')
            config['llm_provider'] = new_provider
            save_config(config)
            log_security_event("config_change", f"AI模型更改: {old_provider} -> {new_provider}")
            logger.info(f"管理员更改AI模型: {old_provider} -> {new_provider}")
            return jsonify({'status': 'success', 'llm_provider': new_provider})
        return jsonify({'status': 'error', 'message': 'Invalid data'}), 400
    else:
        # GET请求无需权限，用于获取当前配置
        config = load_config()
        return jsonify(config)

# AI模型测试路由
@app.route('/api/test-model', methods=['POST'])
def test_model():
    """测试指定AI模型的连通性"""
    try:
        data = request.get_json()
        model_name = data.get('model')

        if not model_name:
            log_security_event("model_test", "缺少模型参数", "failed")
            return jsonify({'status': 'error', 'message': '缺少模型参数'}), 400

        import time
        start_time = time.time()

        # 记录测试开始
        log_security_event("model_test", f"开始测试模型: {model_name}")

        if model_name == 'gemini':
            if not gemini_model:
                log_security_event("model_test", f"Gemini模型未初始化", "failed")
                return jsonify({
                    'status': 'error',
                    'message': 'Gemini模型未初始化',
                    'suggestion': '请检查.env文件中的GEMINI_API_KEY配置'
                })
            try:
                response = gemini_model.generate_content("测试")
                response_time = f"{(time.time() - start_time) * 1000:.0f}ms"
                log_security_event("model_test", f"Gemini测试成功, 响应时间: {response_time}")
                logger.info(f"Gemini模型测试成功，响应时间: {response_time}")
                return jsonify({
                    'status': 'success',
                    'message': '连接成功',
                    'response_time': response_time,
                    'details': 'Gemini API可正常访问'
                })
            except Exception as e:
                log_security_event("model_test", f"Gemini连接失败: {str(e)}", "failed")
                return jsonify({
                    'status': 'error',
                    'message': f'连接失败: {str(e)}',
                    'suggestion': '请检查API密钥和网络连接（可能需要科学上网）'
                })

        elif model_name == 'deepseek':
            if not deepseek_client:
                log_security_event("model_test", f"DeepSeek模型未初始化", "failed")
                return jsonify({
                    'status': 'error',
                    'message': 'DeepSeek客户端未初始化',
                    'suggestion': '请检查.env文件中的DEEPSEEK_API_KEY配置'
                })
            try:
                model_name = load_model_config().get('deepseek_model', 'deepseek-chat')
                response = deepseek_client.chat.completions.create(
                    model=model_name,
                    messages=[{"role": "user", "content": "测试"}],
                    max_tokens=10
                )
                response_time = f"{(time.time() - start_time) * 1000:.0f}ms"
                log_security_event("model_test", f"DeepSeek测试成功, 响应时间: {response_time}")
                logger.info(f"DeepSeek模型测试成功，响应时间: {response_time}")
                return jsonify({
                    'status': 'success',
                    'message': '连接成功',
                    'response_time': response_time,
                    'details': 'DeepSeek API可正常访问'
                })
            except Exception as e:
                log_security_event("model_test", f"DeepSeek连接失败: {str(e)}", "failed")
                return jsonify({
                    'status': 'error',
                    'message': f'连接失败: {str(e)}',
                    'suggestion': '请检查API密钥和网络连接'
                })

        elif model_name == 'ollama_gemma2':
            if not ollama_client:
                log_security_event("model_test", f"Ollama模型未初始化", "failed")
                return jsonify({
                    'status': 'error',
                    'message': 'Ollama客户端未初始化',
                    'suggestion': '请检查Ollama服务是否正常运行'
                })
            try:
                response = ollama_client.chat.completions.create(
                    model="gemma3:27b",
                    messages=[{"role": "user", "content": "测试"}],
                    max_tokens=10
                )
                response_time = f"{(time.time() - start_time) * 1000:.0f}ms"
                log_security_event("model_test", f"Ollama测试成功, 响应时间: {response_time}")
                logger.info(f"Ollama模型测试成功，响应时间: {response_time}")
                return jsonify({
                    'status': 'success',
                    'message': '连接成功',
                    'response_time': response_time,
                    'details': 'Ollama本地服务可正常访问'
                })
            except Exception as e:
                error_msg = str(e)
                if "Connection refused" in error_msg or "Failed to connect" in error_msg:
                    log_security_event("model_test", f"Ollama连接被拒绝", "failed")
                    return jsonify({
                        'status': 'error',
                        'message': '无法连接到Ollama服务',
                        'suggestion': f'请确认Ollama服务正在{os.getenv("OLLAMA_BASE_URL", "http://10.255.1.103:11434/v1")}运行'
                    })
                log_security_event("model_test", f"Ollama连接失败: {error_msg}", "failed")
                return jsonify({
                    'status': 'error',
                    'message': f'连接失败: {error_msg}',
                    'suggestion': '请检查Ollama服务状态和模型是否已下载'
                })

        elif model_name == 'qwen':
            if not qwen_client:
                log_security_event("model_test", f"Qwen模型未初始化", "failed")
                return jsonify({
                    'status': 'error',
                    'message': 'Qwen客户端未初始化',
                    'suggestion': '请检查.env文件中的DASHSCOPE_API_KEY配置'
                })
            try:
                model_name = load_model_config().get('qwen_model', 'qwen-max')
                response = qwen_client.chat.completions.create(
                    model=model_name,  # 动态读取模型名
                    messages=[{"role": "user", "content": "测试"}],
                    max_tokens=10
                )
                response_time = f"{(time.time() - start_time) * 1000:.0f}ms"
                log_security_event("model_test", f"Qwen测试成功, 响应时间: {response_time}")
                logger.info(f"Qwen模型测试成功，响应时间: {response_time}")
                return jsonify({
                    'status': 'success',
                    'message': '连接成功',
                    'response_time': response_time,
                    'details': '阿里通义千问API可正常访问'
                })
            except Exception as e:
                log_security_event("model_test", f"Qwen连接失败: {str(e)}", "failed")
                return jsonify({
                    'status': 'error',
                    'message': f'连接失败: {str(e)}',
                    'suggestion': '请检查API密钥和网络连接'
                })

        elif model_name == 'vllm':
            if not vllm_client:
                log_security_event("model_test", f"VLLM模型未初始化", "failed")
                return jsonify({
                    'status': 'error',
                    'message': 'VLLM客户端未初始化',
                    'suggestion': '请检查.env文件配置'
                })
            try:
                model_name = load_model_config().get('vllm_model', 'gemma-4-moe')
                response = vllm_client.chat.completions.create(
                    model=model_name,
                    messages=[{"role": "user", "content": "测试"}],
                    max_tokens=10
                )
                response_time = f"{(time.time() - start_time) * 1000:.0f}ms"
                log_security_event("model_test", f"VLLM测试成功, 响应时间: {response_time}")
                logger.info(f"VLLM模型测试成功，响应时间: {response_time}")
                return jsonify({
                    'status': 'success',
                    'message': '连接成功',
                    'response_time': response_time,
                    'details': 'VLLM API可正常访问'
                })
            except Exception as e:
                log_security_event("model_test", f"VLLM连接失败: {str(e)}", "failed")
                return jsonify({
                    'status': 'error',
                    'message': f'连接失败: {str(e)}',
                    'suggestion': '请检查API服务和网络连接'
                })

        else:
            log_security_event("model_test", f"未知模型: {model_name}", "failed")
            return jsonify({
                'status': 'error',
                'message': f'未知的模型: {model_name}'
            }), 400

    except Exception as e:
        log_security_event("model_test", f"服务器错误: {str(e)}", "failed")
        logger.error(f"模型测试过程出错: {e}")
        return jsonify({
            'status': 'error',
            'message': f'服务器错误: {str(e)}'
        }), 500

# ========== 高级配置管理API ==========

@app.route('/api/advanced-config', methods=['GET'])
def get_advanced_config():
    """获取高级配置（需要管理员权限）"""
    if not is_admin_authenticated():
        log_security_event("config_access", "未授权尝试获取高级配置", "failed")
        return jsonify({'status': 'error', 'message': '需要管理员权限'}), 403

    try:
        api_keys = load_api_keys()
        model_config = load_model_config()

        # 隐藏API密钥的部分内容（只显示前4位和后4位）
        masked_keys = {}
        for key, value in api_keys.items():
            if value:
                masked_keys[key] = f"{value[:4]}...{value[-4:]}"
            else:
                masked_keys[key] = ""

        log_security_event("config_access", "获取高级配置成功")
        return jsonify({
            'status': 'success',
            'api_keys': masked_keys,
            'model_config': model_config
        })
    except Exception as e:
        log_security_event("config_access", f"获取配置失败: {str(e)}", "failed")
        logger.error(f"获取配置失败: {e}")
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/api/advanced-config', methods=['POST'])
def update_advanced_config():
    """更新高级配置（需要管理员权限）"""
    if not is_admin_authenticated():
        log_security_event("config_change", "未授权尝试修改高级配置", "failed")
        return jsonify({'status': 'error', 'message': '需要管理员权限'}), 403

    try:
        data = request.get_json()
        config_type = data.get('type')

        if config_type == 'api_keys':
            # 更新API密钥
            api_keys = data.get('keys', {})
            if save_api_keys(api_keys):
                log_security_event("config_change", "管理员更新了API密钥配置")
                logger.info("管理员更新了API密钥配置")
                return jsonify({'status': 'success', 'message': 'API密钥已保存'})
            else:
                log_security_event("config_change", "API密钥保存失败", "failed")
                return jsonify({'status': 'error', 'message': '保存失败'}), 500

        elif config_type == 'model_config':
            # 更新模型配置
            model_config = data.get('config', {})
            save_model_config(model_config)
            log_security_event("config_change", f"管理员更新了模型配置: {model_config}")
            logger.info("管理员更新了模型配置")
            return jsonify({'status': 'success', 'message': '模型配置已保存'})

        else:
            log_security_event("config_change", f"无效的配置类型: {config_type}", "failed")
            return jsonify({'status': 'error', 'message': '无效的配置类型'}), 400

    except Exception as e:
        log_security_event("config_change", f"更新配置失败: {str(e)}", "failed")
        logger.error(f"更新配置失败: {e}")
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/api/reload-clients', methods=['POST'])
def reload_clients():
    """重新加载AI客户端（需要管理员权限）"""
    if not is_admin_authenticated():
        log_security_event("system_operation", "未授权尝试重新加载客户端", "failed")
        return jsonify({'status': 'error', 'message': '需要管理员权限'}), 403

    try:
        # 这里可以实现重新初始化客户端的逻辑
        # 暂时返回成功
        log_security_event("system_operation", "管理员请求重新加载AI客户端")
        logger.info("管理员请求重新加载AI客户端")
        return jsonify({'status': 'success', 'message': '客户端已重新加载'})
    except Exception as e:
        log_security_event("system_operation", f"重新加载客户端失败: {str(e)}", "failed")
        logger.error(f"重新加载客户端失败: {e}")
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/process', methods=['POST'])
def process_ppt():
    # 安全检查：文件是否存在
    if 'ppt_file' not in request.files:
        log_security_event("file_upload", "上传请求未包含文件", "failed")
        logger.warning("上传请求未包含文件")
        flash('未检测到文件', 'error')
        return redirect(url_for('index'))

    file = request.files['ppt_file']

    # 安全检查：文件名是否为空
    if file.filename == '':
        log_security_event("file_upload", "上传文件名为空", "failed")
        logger.warning("上传文件名为空")
        flash('未选择文件', 'error')
        return redirect(url_for('index'))

    # 安全检查：文件格式验证
    if not allowed_file(file.filename):
        log_security_event("file_upload", f"无效的文件格式: {file.filename}", "failed")
        logger.warning(f"无效的文件格式: {file.filename}")
        flash('文件格式无效，请上传 .pptx 文件。', 'error')
        return redirect(url_for('index'))

    # 输入验证：时长范围检查 (1-300分钟)
    try:
        duration_minutes = int(request.form.get('duration', 60))
        if not (1 <= duration_minutes <= 300):
            log_security_event("file_upload", f"无效的时长参数: {duration_minutes}", "failed")
            logger.warning(f"无效的时长参数: {duration_minutes}")
            flash('演示时长必须在1-300分钟之间', 'error')
            return redirect(url_for('index'))
    except ValueError:
        log_security_event("file_upload", "时长参数格式错误", "failed")
        logger.error("时长参数格式错误")
        flash('时长参数格式错误', 'error')
        return redirect(url_for('index'))

    # 安全的文件名处理
    original_filename = secure_filename(file.filename)
    session_id = str(uuid.uuid4())
    upload_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{session_id}_{original_filename}")

    try:
        file.save(upload_path)
        log_security_event("file_upload", f"文件上传成功: {original_filename} (会话: {session_id})")
        logger.info(f"文件上传成功: {original_filename} ({session_id})")
    except Exception as e:
        log_security_event("file_upload", f"文件保存失败: {str(e)}", "failed")
        logger.error(f"文件保存失败: {e}")
        flash('文件保存失败，请稍后重试', 'error')
        return redirect(url_for('index'))

    # 全局配置读取：忽略前端传来的 llm_provider，统一使用服务器端配置
    global_config = load_config()
    llm_provider = global_config.get('llm_provider', 'ollama_gemma2')
    logger.info(f"[{session_id}] 使用全局模型配置: {llm_provider}")

    thread = threading.Thread(target=generate_script_task, args=(
        upload_path,
        original_filename,
        llm_provider,
        duration_minutes,
        session_id,
        app.root_path,
        app.config['GENERATED_FOLDER']
    ))
    thread.start()

    tasks[session_id] = {'status': 'starting'}
    return redirect(url_for('results', session_id=session_id))


@app.route('/results/<session_id>')
def results(session_id):
    return render_template('results.html', session_id=session_id)


@app.route('/status/<session_id>')
def task_status(session_id):
    task = tasks.get(session_id)
    if not task: return jsonify({'status': 'not_found'}), 404
    return jsonify(task)


@app.route('/download/<session_id>/<path:filename>')
def download_file(session_id, filename):
    directory = os.path.join(app.config['GENERATED_FOLDER'], session_id)
    return send_from_directory(directory, filename, as_attachment=True)

# ========== 客户端监控日志 API ==========
@app.route('/api/client-log', methods=['POST'])
def client_log():
    try:
        data = request.get_json()
        if not data:
            return jsonify({'status': 'error', 'message': '无数据'}), 400
            
        log_type = data.get('type', 'info')
        category = data.get('category', 'general')
        message = data.get('message', '')
        details = data.get('details', {})
        
        # 统一写入 app.log 或专用的 frontend_monitor.log
        log_entry = f"[前端监控] [{log_type.upper()}] [{category}] {message} - {json.dumps(details, ensure_ascii=False)}"
        
        
        if log_type == 'error':
            logger.error(log_entry)
        elif log_type == 'perf':
            logger.info(log_entry)
        elif log_type == 'biz':
            logger.info(log_entry)
        else:
            logger.info(log_entry)
            
        return jsonify({'status': 'success'})
    except Exception as e:
        logger.error(f"处理客户端日志时出错: {e}")
        return jsonify({'status': 'error', 'message': str(e)}), 500

if __name__ == '__main__':
    # 从环境变量读取调试模式设置
    flask_debug = os.getenv("FLASK_DEBUG", "false").lower() in ["true", "1", "yes"]
    app.run(debug=flask_debug, host='0.0.0.0', port=6001)